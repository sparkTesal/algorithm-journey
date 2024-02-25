import os
from pptx import Presentation
from PIL import Image


def ppt_to_images(ppt_path, images_path):
    presentation = Presentation(ppt_path)
    img_list = []

    # 遍历每一个幻灯片
    for i, slide in enumerate(presentation.slides):
        # 将每个幻灯片保存为一个临时图片文件
        img_path = f'temp_slide_{i}.jpg'
        slide.shapes._spTree.export(img_path, format='jpg')
        img_list.append(Image.open(img_path))

    # 获取所有幻灯片图片的最大宽度和总高度
    widths, heights = zip(*(i.size for i in img_list))
    max_width = max(widths)
    total_height = sum(heights)

    # 创建一个足够容纳所有幻灯片的长图片
    long_img = Image.new('RGB', (max_width, total_height))

    # 拼接图片
    y_offset = 0
    for img in img_list:
        long_img.paste(img, (0, y_offset))
        y_offset += img.size[1]
        img.close()  # 关闭图片以释放资源

    # 保存最终的长图片
    long_img.save(images_path)
    long_img.close()

    # 删除临时图片文件
    for i in range(len(img_list)):
        os.remove(f'temp_slide_{i}.jpg')


def convert_folder_to_images(folder_path):
    for file in os.listdir(folder_path):
        if file.endswith(('.ppt', '.pptx')):
            ppt_path = os.path.join(folder_path, file)
            images_path = os.path.join(folder_path, f'{os.path.splitext(file)[0]}.jpg')
            ppt_to_images(ppt_path, images_path)
            print(f'Converted {file} to image.')


if __name__ == '__main__':
    # 转换当前文件夹下的所有PPT文件
    convert_folder_to_images('ppt')
